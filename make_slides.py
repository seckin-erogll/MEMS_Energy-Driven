"""
make_slides.py — Generate a simple paper-outline slideshow for the MEMS PINN work.
Run: python make_slides.py
Output: MEMS_PINN_slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── constants ─────────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

BLACK  = RGBColor(0x10, 0x10, 0x10)
ACCENT = RGBColor(0x1A, 0x5C, 0xB0)   # steel blue — used sparingly
GRAY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xFA)   # very pale blue for header bar

# ── helpers ───────────────────────────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()   # no border
    return shape


def header_bar(slide, title_text, subtitle=None):
    """Top bar: accent rectangle + white title text."""
    bar_h = Inches(1.05)
    add_rect(slide, 0, 0, W, bar_h, ACCENT)
    add_textbox(slide, title_text,
                left=Inches(0.35), top=Inches(0.08),
                width=Inches(11.5), height=bar_h,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    left=Inches(0.35), top=Inches(0.65),
                    width=Inches(12), height=Inches(0.45),
                    size=14, bold=False, color=WHITE, italic=True)


def bullet_block(slide, items, left, top, width, height,
                 base_size=16, indent_size=13, gap=Inches(0.0)):
    """
    items: list of (indent_level, text)  indent 0=main, 1=sub
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for level, text in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        p.space_before = Pt(4 if level == 0 else 1)
        run = p.add_run()

        if level == 0:
            run.text = "▸  " + text
            run.font.size  = Pt(base_size)
            run.font.bold  = True
            run.font.color.rgb = BLACK
        else:
            run.text = "     –  " + text
            run.font.size  = Pt(indent_size)
            run.font.bold  = False
            run.font.color.rgb = GRAY

    return txb


def two_col(slide, left_items, right_items,
            top=Inches(1.2), height=Inches(5.8),
            base_size=16, indent_size=13):
    col_w = Inches(6.2)
    bullet_block(slide, left_items,
                 left=Inches(0.35), top=top,
                 width=col_w, height=height,
                 base_size=base_size, indent_size=indent_size)
    bullet_block(slide, right_items,
                 left=Inches(6.9), top=top,
                 width=col_w, height=height,
                 base_size=base_size, indent_size=indent_size)


def code_box(slide, code_text, left, top, width, height, size=13):
    add_rect(slide, left, top, width, height, RGBColor(0xF5, 0xF5, 0xF5))
    txb = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.08),
        width - Inches(0.24), height - Inches(0.16)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x5E)   # dark navy monospace feel


def slide_number(slide, n, total=11):
    add_textbox(slide, f"{n} / {total}",
                left=Inches(12.3), top=Inches(7.1),
                width=Inches(0.9), height=Inches(0.35),
                size=11, color=GRAY, align=PP_ALIGN.RIGHT)


# ── build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # blank layout for all slides


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, Inches(2.5), W, Inches(2.6), ACCENT)

add_textbox(s,
    "Energy-Driven Physics-Informed Neural Network\n"
    "for Multilayer MEMS Diaphragm Pressure Sensors",
    left=Inches(0.7), top=Inches(2.6),
    width=Inches(11.9), height=Inches(1.8),
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s,
    "Föppl–von Kármán Variational Formulation  |  "
    "Classical Laminated Plate Theory  |  "
    "Hard Boundary Conditions  |  Contact Enforcement",
    left=Inches(0.7), top=Inches(4.35),
    width=Inches(11.9), height=Inches(0.55),
    size=14, bold=False, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

add_textbox(s,
    "Parylene-C / Gold / Parylene-C  ·  Circular Clamped Diaphragm  ·  Capacitive Sensor",
    left=Inches(0.7), top=Inches(5.3),
    width=Inches(11.9), height=Inches(0.5),
    size=16, color=ACCENT, align=PP_ALIGN.CENTER)

slide_number(s, 1)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation & Problem Statement
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "1.  Motivation & Problem Statement")
slide_number(s, 2)

two_col(s,
    left_items=[
        (0, "Device"),
        (1, "Clamped circular diaphragm: parylene-C / Au / parylene-C"),
        (1, "Radius a = 150–500 µm,  air gap ag = 5–15 µm"),
        (1, "Uniform pressure P = 10 Pa – 20 kPa"),
        (0, "Goal"),
        (1, "Predict full deflection profile w(r) for any geometry & pressure"),
        (1, "Without re-running COMSOL for each new design"),
        (0, "Challenge: large deflections"),
        (1, "w/h can reach 5–10×  →  linear plate theory fails"),
        (1, "Membrane stiffening must be captured"),
        (1, "Contact with electrode at w = −ag  must be enforced"),
    ],
    right_items=[
        (0, "Why not just FEA?"),
        (1, "COMSOL: minutes–hours per geometry point"),
        (1, "Design optimisation needs 10³–10⁴ evaluations"),
        (1, "Surrogate model gives <1 ms inference"),
        (0, "Why not Kirchhoff (Atik analytical)?"),
        (1, "w(r) = −Pa⁴/(64D*) · (1−ξ²)²"),
        (1, "Valid only for w << thickness"),
        (1, "Ignores membrane stiffening  →  overestimates deflection"),
        (1, "No contact enforcement"),
        (0, "Approach"),
        (1, "PINN trained by minimising the physical energy functional"),
        (1, "No labelled data required during main training"),
    ]
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Physics: Föppl–von Kármán
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "2.  Governing Physics: Föppl–von Kármán (FvK)")
slide_number(s, 3)

bullet_block(s,
    items=[
        (0, "Kirchhoff (linear) plate theory — valid only for w << h"),
        (1, "Bending stiffness D* resists deflection"),
        (1, "Equilibrium: D∇⁴w + P = 0"),
        (0, "FvK extension — large deflections"),
        (1, "Plate also stretches in-plane as it bends"),
        (1, "In-plane tension (membrane stiffening) increases effective stiffness"),
        (1, "Two coupled unknowns: w(r) transverse,  u(r) radial in-plane"),
        (0, "Nonlinear membrane strains  (the key physics)"),
        (1, "εᵣ = u'(r) + ½(w'(r))²      ← ½(w')² is the geometric nonlinearity"),
        (1, "εθ = u(r) / r"),
        (1, "Dropping u or the ½(w')² term loses all membrane stiffening"),
        (0, "No closed-form solution for general multilayer geometry  →  need a surrogate"),
    ],
    left=Inches(0.35), top=Inches(1.15),
    width=Inches(8.5), height=Inches(6.0),
)

code_box(s,
    "Bending energy density:\n"
    "  ½ D* (w'' + w'/r)²\n"
    "\n"
    "Membrane energy density:\n"
    "  ½ [A₁₁εᵣ² + 2A₁₂εᵣεθ + A₁₁εθ²]\n"
    "\n"
    "Pressure work density:\n"
    "  P · w    (w < 0  →  decreases Π)",
    left=Inches(9.1), top=Inches(1.3),
    width=Inches(3.95), height=Inches(3.0),
    size=13,
)

add_textbox(s, "D*, A₁₁, A₁₂ from Classical\nLaminated Plate Theory\n(computed analytically)",
            left=Inches(9.1), top=Inches(4.5),
            width=Inches(3.95), height=Inches(1.5),
            size=13, color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Classical Laminated Plate Theory
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "3.  Material Model: Classical Laminated Plate Theory (CLT)")
slide_number(s, 4)

two_col(s,
    left_items=[
        (0, "3-layer stack: parylene-C / gold / parylene-C"),
        (1, "Parylene-C:  E = 3.2 GPa,  ν = 0.33"),
        (1, "Gold:        E = 70  GPa,  ν = 0.44  (22× stiffer)"),
        (0, "Neutral axis z_NA"),
        (1, "z_NA = Σ Ēₖ tₖ z_{c,k}  /  Σ Ēₖ tₖ"),
        (1, "Gold is stiff and not at geometric midplane"),
        (1, "Pulls z_NA upward  →  asymmetric bending stiffness"),
        (0, "Effective bending stiffness D*"),
        (1, "D* = (1/3) Σ Ēₖ (bₖ³ − bₖ₋₁³)"),
        (1, "bₖ = signed distance from z_NA to layer top"),
        (1, "Cubic: thin gold far from NA contributes like I-beam flange"),
        (0, "Membrane stiffness A₁₁, A₁₂"),
        (1, "A₁₁ = Σ Ēₖ tₖ   (axial)"),
        (1, "A₁₂ = Σ νₖ Ēₖ tₖ  (Poisson coupling)"),
    ],
    right_items=[
        (0, "Why compute CLT analytically, not learn it?"),
        (1, "D*, A₁₁, A₁₂ are exact — no approximation needed"),
        (1, "Embedding exact physics reduces what the network must learn"),
        (1, "Generalises outside training range without extrapolation error"),
        (0, "Nominal values  (t1=1, t2=0.2, t3=4 µm)"),
        (1, "z_NA = 1.89 µm  (geometric midplane = 2.6 µm)"),
        (1, "D* = 6.19 × 10⁻⁸ N·m"),
        (1, "A₁₁ = 3.53 × 10⁴ N/m"),
        (1, "ν_eff = A₁₂/A₁₁ = 0.384"),
        (0, "Computed fresh for every training batch"),
        (1, "Geometry (t1, t2, t3) is sampled stochastically each step"),
        (1, "D*, A₁₁, A₁₂ evaluated analytically at batch geometry"),
    ],
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Variational Energy Formulation
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "4.  Training Signal: Variational Energy Functional")
slide_number(s, 5)

bullet_block(s,
    items=[
        (0, "Total potential energy"),
        (1, "Π[w, u] = ∫₀ᵃ 2πr [ ½D*(w''+w'/r)²  +  ½(A₁₁εᵣ² + 2A₁₂εᵣεθ + A₁₁εθ²)  +  Pw ] dr"),
        (0, "Why energy, not PDE residual?"),
        (1, "Energy is a single scalar  —  simpler gradient signal than pointwise PDE enforcement"),
        (1, "At the minimum, δΠ/δw = 0 recovers D∇⁴w + P = 0  (FvK Euler–Lagrange)"),
        (1, "Natural boundary conditions satisfied automatically at the minimum"),
        (1, "No labelled data needed — physics alone constrains the solution"),
        (0, "Why +Pw (not −Pw)?"),
        (1, "w < 0 for downward deflection  →  Pw < 0  →  decreases Π as plate deflects"),
        (1, "Energy minimisation drives the plate downward — physically correct"),
        (1, "Wrong sign locks the network at w = 0 (deflection is energetically unfavourable)"),
        (0, "Numerical integration: 32-node Gauss–Legendre quadrature on ξ ∈ (0, 1)"),
        (1, "Differentiable weighted sum  →  backpropagation flows through the integral"),
        (1, "ξ = 0 excluded  —  avoids 1/r singularity in εθ = u/r"),
        (1, "Exponentially accurate for smooth (tanh) networks — 32 nodes is sufficient"),
        (0, "Per-sample energy normalisation  (critical for multi-scale training)"),
        (1, "Energy at P = 20 kPa  ~  10⁻⁸ J  vs  P = 50 Pa  ~  10⁻¹² J  →  4 orders of magnitude"),
        (1, "Π_norm = Π_i / |Π_i|  →  every sample contributes ±1 to the gradient"),
        (1, "Without this: P = 100 Pa error = 232%  |  with this: 2.2%"),
    ],
    left=Inches(0.35), top=Inches(1.12),
    width=Inches(12.6), height=Inches(6.1),
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Network Architecture
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "5.  Network Architecture")
slide_number(s, 6)

bullet_block(s,
    items=[
        (0, "Input: 7 normalised features  [ξ, P̂, t̂₁, t̂₂, t̂₃, â, âg]"),
        (1, "ξ = r/a  (radial coordinate, already in [0,1] — network learns shapes, not sizes)"),
        (1, "P̂ = log(1 + P/1000) / log(21)  —  log-scale: 10 Pa–20 kPa span 3 decades"),
        (1, "t̂₁, t̂₂, t̂₃, â, âg  —  min-max to [0,1] within training geometry range"),
        (0, "Trunk: 6 fully-connected layers, 128 neurons, Tanh activation throughout"),
        (1, "Why Tanh?  Bending energy needs w''  (second autograd derivative through the network)"),
        (1, "ReLU  →  w'' = 0 almost everywhere  →  no bending energy gradient"),
        (1, "GELU  →  w'' noisy and oscillatory  →  unstable energy computation"),
        (1, "Tanh  →  smooth, analytic, well-behaved second derivatives  →  stable training"),
        (0, "Two independent output heads (Linear 128→1 each)"),
        (1, "Head w  →  transverse deflection w(r)  [the observable]"),
        (1, "Head u  →  radial in-plane displacement u(r)  [needed for εᵣ = u' + ½(w')²]"),
        (1, "Without u: εᵣ = ½(w')² only  →  misses actual in-plane displacement  →  wrong stiffness"),
        (0, "~100 k parameters total  |  float32  |  trained on GPU"),
    ],
    left=Inches(0.35), top=Inches(1.12),
    width=Inches(8.6), height=Inches(6.1),
    base_size=15, indent_size=12,
)

code_box(s,
    "Input  [7]\n"
    "  ↓  Linear(7→128) + Tanh\n"
    "  ↓  Linear(128→128) + Tanh  ×5\n"
    "  ↓\n"
    "  ┌──────────┬──────────┐\n"
    "  head_w    head_u\n"
    "  Lin(128→1) Lin(128→1)\n"
    "  ↓            ↓\n"
    "  hard BC    hard BC\n"
    "  ↓            ↓\n"
    " w(r) [m]   u(r) [m]",
    left=Inches(9.15), top=Inches(1.25),
    width=Inches(3.85), height=Inches(3.3),
    size=13,
)

add_textbox(s,
    "Xavier uniform init\nZero-init output heads",
    left=Inches(9.15), top=Inches(4.7),
    width=Inches(3.85), height=Inches(0.6),
    size=12, color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Hard BCs & Contact Constraint
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "6.  Hard Boundary Conditions & Contact Constraint")
slide_number(s, 7)

two_col(s,
    left_items=[
        (0, "Clamped rim: algebraic enforcement"),
        (1, "w_profile = (1 − ξ²)² · w_raw · W_ref"),
        (1, "(1−ξ²)² is the exact Kirchhoff clamped-plate polynomial"),
        (1, "Guarantees: w(ξ=1) = 0,  w'(ξ=1) = 0,  symmetry at ξ = 0"),
        (1, "Whatever the network outputs, BCs are mathematically exact"),
        (0, "In-plane BC"),
        (1, "u = ξ(1−ξ) · u_raw · U_ref"),
        (1, "Guarantees: u(0) = 0 (symmetry),  u(1) = 0 (clamped rim)"),
        (0, "Why hard BCs instead of soft BC penalty?"),
        (1, "Soft penalty: BCs can still be violated; extra hyperparameter"),
        (1, "Hard BC: zero gradient budget wasted on BC enforcement"),
        (1, "All network capacity goes toward learning interior physics"),
        (0, "Physical output scaling: W_ref = ag_phys"),
        (1, "Maximum deflection is −ag  →  w_raw lives in O(1)"),
        (1, "Consistent across all geometries regardless of ag"),
    ],
    right_items=[
        (0, "Contact constraint: w ≥ −ag"),
        (1, "Smooth-max: w = ½(w_p + (−ag) + √((w_p−(−ag))² + ε²))"),
        (1, "ε = 50 nm  →  transition width at contact line"),
        (0, "Why not torch.clamp(w, min=−ag)?"),
        (1, "clamp: gradient = 0 in contact zone"),
        (1, "Network receives no learning signal where plate touches electrode"),
        (1, "Smooth-max: gradient nonzero everywhere  →  learns contact physics"),
        (0, "Why not κ-obstacle penalty?"),
        (1, "Contact penalty ~10⁻¹² J  vs  pressure work ~10⁻⁸ J at 15 kPa"),
        (1, "Ratio 1:10,000  →  constraint dominated by pressure term"),
        (1, "Even at κ = 10¹⁰: still 1.9 µm penetration"),
        (0, "Smooth-max result"),
        (1, "Hard mathematical guarantee: w ≥ −ag to within 50 nm"),
        (1, "0 / 200 penetration violations across random geometry sweep"),
    ],
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Training
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "7.  Training Strategy")
slide_number(s, 8)

two_col(s,
    left_items=[
        (0, "Geometry sampling  (each step)"),
        (1, "64 geometries × 8 pressures = 512 samples per step"),
        (1, "t1, t3 ~ Uniform[1–5 µm],  t2 ~ [0.2–0.5 µm]"),
        (1, "a ~ Uniform[150–500 µm],  ag ~ [5–15 µm]"),
        (1, "P ~ Log-Uniform[10 Pa – 20 kPa]"),
        (0, "COMSOL warmup (steps 0–5000, λ annealing 1→0)"),
        (1, "Anchors network in correct deflection regime"),
        (1, "Prevents spurious minima (flat plate, upward deflection)"),
        (1, "L_total = L_energy + λ · L_COMSOL_MSE"),
        (1, "After step 5000: λ = 0, pure physics only"),
        (0, "Optimiser: Adam"),
        (1, "LR: 1×10⁻³ (warmup) → cosine decay to 1×10⁻⁴"),
        (1, "Gradient clipping: max_norm = 1.0"),
        (1, "50 000 steps total  |  ~55 min on GPU"),
    ],
    right_items=[
        (0, "Why stochastic geometry sampling?"),
        (1, "Network must generalise over 5D geometry space"),
        (1, "Random sampling covers the space without grid explosion"),
        (1, "5D uniform grid at 10 pts/dim = 10⁵ points  →  impractical"),
        (0, "Why log-uniform pressure sampling?"),
        (1, "Three-decade range: 10 Pa – 20 kPa"),
        (1, "Linear sampling starves low-P regime"),
        (1, "Log-uniform: equal samples per decade"),
        (0, "Why gradient clipping?"),
        (1, "Energy gradient can spike at initialisation (random w far from equilibrium)"),
        (1, "Clipping prevents early instability without affecting steady-state"),
        (0, "Loss history"),
        (1, "L_energy converges from ~+1 (random init) to ~−0.7"),
        (1, "Negative energy: pressure work dominates  →  physically correct"),
    ],
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Deflection Profile Results
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "8.  Results: Deflection Profiles — 5 Geometries")
slide_number(s, 9)

bullet_block(s,
    items=[
        (0, "Three-way comparison: PINN (solid)  |  COMSOL FEA (dashed)  |  Atik/Kirchhoff (dotted)"),
        (0, "5 geometries spanning the design space"),
        (1, "Small thin       a=227 µm,  t1/t3=1.6/1.1 µm,  ag=6.0 µm"),
        (1, "Medium balanced  a=290 µm,  t1/t3=1.7/2.1 µm,  ag=7.9 µm"),
        (1, "Large thick      a=460 µm,  t1/t3=4.3/3.9 µm,  ag=11.1 µm"),
        (1, "Small large gap  a=179 µm,  t1/t3=1.7/2.1 µm,  ag=13.1 µm"),
        (1, "Large thin       a=462 µm,  t1/t3=1.6/1.7 µm,  ag=5.4 µm"),
        (0, "Key observations"),
        (1, "Low pressure (500 Pa): all three methods agree  →  linear regime"),
        (1, "High pressure (10–20 kPa): Atik overestimates deflection (dotted below solid/dashed)"),
        (1, "Membrane stiffening captured by PINN/COMSOL  →  less deflection than Kirchhoff predicts"),
        (1, "Touch-mode: PINN & COMSOL both plateau at w = −ag;  Atik is clipped artificially"),
        (1, "Atik predicts touch-down at lower pressure than reality  (ignores membrane stiffening)"),
        (0, "PINN vs COMSOL agreement is the key validation"),
        (1, "Solid and dashed nearly overlapping across all pressures and geometries"),
    ],
    left=Inches(0.35), top=Inches(1.12),
    width=Inches(12.6), height=Inches(6.1),
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Quantitative Validation
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "9.  Quantitative Validation — 2144 COMSOL Files")
slide_number(s, 10)

two_col(s,
    left_items=[
        (0, "Validation set"),
        (1, "2144 COMSOL deflection profiles"),
        (1, "Geometries span the full 5D training box"),
        (1, "8 pressures per geometry (500 Pa → 20 kPa)"),
        (0, "RMSE: PINN vs COMSOL  (µm)"),
        (1, "Median RMSE:   0.15 µm"),
        (1, "Mean RMSE:     0.22 µm"),
        (1, "p90 RMSE:      0.51 µm"),
        (1, "Max RMSE:      1.10 µm  (worst: a ≈ 492 µm, boundary)"),
        (0, "Linear regime accuracy"),
        (1, "P = 100 Pa Kirchhoff error (w_PINN vs w_exact): 2.2%"),
        (1, "Confirms energy minimum recovers correct PDE solution"),
        (0, "Contact enforcement"),
        (1, "0 / 200 penetration violations in random geometry sweep"),
        (1, "Hard smooth-max constraint holds across full pressure range"),
    ],
    right_items=[
        (0, "Error breakdown by pressure"),
        (1, "Low P (≤ 1 kPa): RMSE < 0.05 µm  (linear regime, easy)"),
        (1, "Mid P (1–5 kPa): RMSE 0.1–0.3 µm  (membrane stiffening onset)"),
        (1, "High P (≥ 10 kPa): RMSE 0.3–1.1 µm  (touch mode, hardest)"),
        (0, "What drives the remaining error?"),
        (1, "Large-plate geometries (a > 400 µm) show highest RMSE"),
        (1, "Near training-box boundary  →  less training coverage"),
        (1, "Touch-mode transition radius prediction (contact area boundary)"),
        (0, "Context"),
        (1, "Air gap range: 5–15 µm  →  RMSE / ag = 1–7%"),
        (1, "Comparable or better than other PINN plate surrogates"),
        (1, "Inference: < 1 ms  vs  COMSOL: minutes per point"),
    ],
    base_size=15, indent_size=12,
)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusions
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "10.  Conclusions & Design Decisions")
slide_number(s, 11)

two_col(s,
    left_items=[
        (0, "What was built"),
        (1, "PINN surrogate for multilayer MEMS diaphragm over 5D geometry space"),
        (1, "Trained purely on physics energy (no FEA labels in main training)"),
        (1, "Median RMSE 0.15 µm vs COMSOL,  0 contact violations"),
        (0, "Key design decisions (and why)"),
        (1, "Tanh activations  →  clean autograd second derivatives for bending energy"),
        (1, "Two heads (w, u)  →  correct FvK membrane strain εᵣ = u' + ½(w')²"),
        (1, "Hard BCs (1−ξ²)²  →  exact enforcement, no wasted gradient budget"),
        (1, "Smooth-max contact  →  hard guarantee + nonzero gradient in touch mode"),
        (1, "Per-sample energy normalisation  →  equal learning across 4-decade pressure range"),
        (1, "Log-scale pressure input  →  equal representation of all pressure decades"),
        (1, "W_ref = ag  →  O(1) network outputs across all air-gap geometries"),
    ],
    right_items=[
        (0, "What the figures show"),
        (1, "PINN ≈ COMSOL at all pressures and geometries  →  surrogate is valid"),
        (1, "Kirchhoff overestimates deflection at P > 1 kPa  →  FvK is necessary"),
        (1, "Touch-mode correctly captured; contact radius matches COMSOL"),
        (0, "Limitations"),
        (1, "Max RMSE 1.10 µm for large plates (a > 400 µm)"),
        (1, "Touch-mode transition slightly imprecise at very high P"),
        (0, "Possible extensions"),
        (1, "Wider network (256-dim) or 100k steps to close large-plate gap"),
        (1, "Capacitance C(P) prediction directly from w(r) (Eroglu 2025)"),
        (1, "Extend to elliptical or rectangular diaphragm geometries"),
        (1, "Uncertainty quantification via ensemble or dropout"),
    ],
    base_size=15, indent_size=12,
)


# ── save ──────────────────────────────────────────────────────────────────────
out = "MEMS_PINN_slides.pptx"
prs.save(out)
print(f"Saved → {out}")
